#!/usr/bin/env python3
"""
build_deck.py — Teroxx Portfolio Allocation App: Handover & Self-Hosting deck.

Native, editable PowerPoint (python-pptx) in the official Teroxx identity
("blue-hour" palette, Söhne / Sometimes Times). Convert to PDF afterwards with
LibreOffice:  soffice --headless --convert-to pdf <deck>.pptx

Audience: Research (Jannick Broering, Leonardo Larieira — non-technical) and
IT (Emil Jorgensen — technical). One deck, layered: a plain-language front half,
a technical + migration-runbook back half.
"""
from __future__ import annotations

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

HERE = os.path.dirname(os.path.abspath(__file__))
IMG = os.path.join(os.path.dirname(HERE), "app", "static", "img")
LOGO_WHITE = os.path.join(IMG, "logo-white.png")
LOGO_DARK = os.path.join(IMG, "logo-dark.png")

# ── Teroxx "blue-hour" palette ────────────────────────────────────────────
NIGHTBLUE   = RGBColor(0x01, 0x06, 0x26)
DEEP_INDIGO = RGBColor(0x06, 0x0D, 0x43)
HEADER_BG   = RGBColor(0x17, 0x1E, 0x4B)
ELECTRIC    = RGBColor(0x0B, 0x68, 0x8C)
SANDSTONE   = RGBColor(0xBF, 0xB3, 0xA8)
CREAM       = RGBColor(0xEC, 0xE8, 0xE5)
CREAM_DARK  = RGBColor(0xD9, 0xD2, 0xCB)
EMBER       = RGBColor(0xD0, 0x66, 0x43)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
SUCCESS     = RGBColor(0x1A, 0x8A, 0x4A)
DANGER      = RGBColor(0xC0, 0x43, 0x2A)
INK         = RGBColor(0x0E, 0x14, 0x30)
INK_70      = RGBColor(0x4A, 0x51, 0x70)
INK_40      = RGBColor(0x90, 0x95, 0xAB)
GRID        = RGBColor(0xC9, 0xCC, 0xD9)
PANEL       = RGBColor(0xF4, 0xF2, 0xEF)   # warm off-white panel

SANS      = "Söhne"
SANS_MED  = "Söhne Kräftig"
SANS_SEMI = "Söhne Halbfett"
SANS_BOLD = "Söhne Fett"
SERIF     = "Sometimes Times Medium"

# audience accent chips
RESEARCH_C = ELECTRIC
IT_C       = EMBER

EMU_W, EMU_H = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width = EMU_W
prs.slide_height = EMU_H
BLANK = prs.slide_layouts[6]


# ── low-level helpers ──────────────────────────────────────────────────────
def slide():
    return prs.slides.add_slide(BLANK)


def rect(s, l, t, w, h, fill=None, line=None, line_w=None, shadow=False):
    sp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = line_w or Pt(1)
    sp.shadow.inherit = False
    if shadow:
        el = sp._element.spPr
        # subtle drop shadow
        ns = ('<a:effectLst xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
              '<a:outerShdw blurRad="40000" dist="20000" dir="5400000" rotWithShape="0">'
              '<a:srgbClr val="0E1430"><a:alpha val="16000"/></a:srgbClr></a:outerShdw></a:effectLst>')
        from pptx.oxml import parse_xml
        el.append(parse_xml(ns))
    return sp


def line_shape(s, l, t, w, h, color, weight=Pt(1)):
    sp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    sp.line.fill.background(); sp.shadow.inherit = False
    return sp


def _set_run(r, text, font, size, color, bold=False, italic=False, spacing=None):
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.italic = italic
    # ensure east-asian/complex also map to same font
    rPr = r._r.get_or_add_rPr()
    for tag in ("latin", "cs"):
        e = rPr.find(qn("a:" + tag))
        if e is None:
            e = rPr.makeelement(qn("a:" + tag), {}); rPr.append(e)
        e.set("typeface", font)
    if spacing is not None:
        rPr.set("spc", str(int(spacing * 100)))


def text(s, l, t, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         line_spacing=1.0, space_after=0, wrap=True):
    """runs: list of paragraphs; each paragraph is a list of (text,font,size,color,bold,italic) tuples."""
    tb = s.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        for spec in para:
            txt, font, size, color = spec[0], spec[1], spec[2], spec[3]
            bold = spec[4] if len(spec) > 4 else False
            italic = spec[5] if len(spec) > 5 else False
            spc = spec[6] if len(spec) > 6 else None
            r = p.add_run()
            _set_run(r, txt, font, size, color, bold, italic, spc)
    return tb


def bg(s, color):
    rect(s, 0, 0, EMU_W, EMU_H, fill=color)


# ── page furniture ─────────────────────────────────────────────────────────
_PAGE = {"n": 0}


def content_header(s, kicker, title, chip=None):
    """White content slide header: kicker + serif title + rule."""
    bg(s, WHITE)
    rect(s, 0, 0, EMU_W, Inches(1.28), fill=WHITE)
    text(s, Inches(0.6), Inches(0.34), Inches(9.5), Inches(0.3),
         [[(kicker.upper(), SANS_SEMI, 10.5, ELECTRIC, True, False, 2.2)]])
    text(s, Inches(0.6), Inches(0.58), Inches(11.4), Inches(0.7),
         [[(title, SERIF, 25, DEEP_INDIGO)]])
    line_shape(s, Inches(0.6), Inches(1.24), Inches(12.13), Pt(2), ELECTRIC)
    if chip:
        audience_chip(s, chip, Inches(11.0), Inches(0.38))
    footer(s)


def audience_chip(s, kind, l, t):
    """kind: 'research', 'it', 'both'"""
    label, col = {
        "research": ("FOR RESEARCH", RESEARCH_C),
        "it": ("FOR IT", IT_C),
        "both": ("RESEARCH + IT", DEEP_INDIGO),
    }[kind]
    w = Inches(1.75)
    r = rect(s, l, t, w, Inches(0.34), fill=None, line=col, line_w=Pt(1.25))
    text(s, l, t, w, Inches(0.34),
         [[(label, SANS_SEMI, 8.5, col, True, False, 1.5)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def footer(s, dark=False):
    _PAGE["n"] += 1
    col = CREAM_DARK if dark else INK_40
    text(s, Inches(0.6), Inches(7.06), Inches(9.0), Inches(0.3),
         [[("Teroxx  ·  Portfolio Allocation App — Handover  ·  Confidential",
            SANS, 8, col)]], anchor=MSO_ANCHOR.MIDDLE)
    text(s, Inches(11.6), Inches(7.06), Inches(1.13), Inches(0.3),
         [[(str(_PAGE["n"]), SANS_SEMI, 8, col, True)]],
         align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)


def card(s, l, t, w, h, fill=PANEL, line=GRID):
    return rect(s, l, t, w, h, fill=fill, line=line, line_w=Pt(0.75))


def bullets(s, l, t, w, h, items, size=13, color=INK, gap=8, marker="—",
            marker_color=None, lh=1.12, bold_lead=False):
    """items: list of strings, or (lead, rest) tuples where lead is bolded."""
    tb = s.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    mc = marker_color or ELECTRIC
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = lh; p.space_after = Pt(gap); p.space_before = Pt(0)
        rm = p.add_run(); _set_run(rm, marker + "  ", SANS_SEMI, size, mc, True)
        if isinstance(it, tuple):
            lead, rest = it
            r1 = p.add_run(); _set_run(r1, lead, SANS_SEMI, size, color, True)
            if rest:
                r2 = p.add_run(); _set_run(r2, rest, SANS, size, INK_70)
        else:
            r = p.add_run(); _set_run(r, it, SANS, size, color)
    return tb


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 1 — COVER
# ═══════════════════════════════════════════════════════════════════════════
def s_cover():
    s = slide(); bg(s, NIGHTBLUE)
    # deep indigo band + ember accent
    rect(s, 0, Inches(4.55), EMU_W, Inches(2.95), fill=DEEP_INDIGO)
    line_shape(s, Inches(0.9), Inches(2.35), Inches(1.7), Pt(3), EMBER)
    if os.path.exists(LOGO_WHITE):
        s.shapes.add_picture(LOGO_WHITE, Inches(0.85), Inches(0.7), height=Inches(0.52))
    text(s, Inches(0.9), Inches(1.5), Inches(11), Inches(0.4),
         [[("PROJECT HANDOVER  ·  JULY 2026", SANS_SEMI, 12, SANDSTONE, True, False, 3.0)]])
    text(s, Inches(0.9), Inches(2.55), Inches(11.6), Inches(1.9),
         [[("Portfolio Allocation App", SERIF, 46, WHITE)],
          [("Status, architecture & how to run it on Teroxx cloud", SERIF, 22, CREAM)]],
         line_spacing=1.02)
    text(s, Inches(0.9), Inches(4.95), Inches(11), Inches(0.4),
         [[("Prepared by Aleksander Biesaga  —  Digital Assets, Teroxx", SANS_MED, 13.5, WHITE)]])
    text(s, Inches(0.9), Inches(5.5), Inches(11), Inches(1.3),
         [[("Handover to", SANS_SEMI, 10.5, SANDSTONE, True, False, 2.0)],
          [("Research — Jannick Broering, Leonardo Larieira", SANS, 13, CREAM)],
          [("IT — Emil Jorgensen", SANS, 13, CREAM)]],
         line_spacing=1.25, space_after=2)
    text(s, Inches(0.9), Inches(6.9), Inches(11), Inches(0.3),
         [[("Confidential — internal use only", SANS, 9, INK_40)]])


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 2 — HOW TO READ THIS DECK / CONTENTS
# ═══════════════════════════════════════════════════════════════════════════
def s_contents():
    s = slide(); content_header(s, "Orientation", "How to read this deck", chip="both")
    text(s, Inches(0.6), Inches(1.5), Inches(12.1), Inches(0.7),
         [[("This one deck serves two audiences. ", SANS_SEMI, 13, INK, True),
           ("Every slide is tagged so you can read only what is relevant to you — "
            "or read it all for the full picture.", SANS, 13, INK_70)]],
         line_spacing=1.15)
    # two audience cards
    card(s, Inches(0.6), Inches(2.3), Inches(5.9), Inches(1.55), fill=PANEL, line=RESEARCH_C)
    line_shape(s, Inches(0.6), Inches(2.3), Inches(0.09), Inches(1.55), RESEARCH_C)
    text(s, Inches(0.85), Inches(2.5), Inches(5.5), Inches(1.2),
         [[("FOR RESEARCH", SANS_SEMI, 11, RESEARCH_C, True, False, 1.5)],
          [("What the app does, the day-to-day workflow, the model, and where it "
            "lives. No code required.", SANS, 12, INK_70)]],
         line_spacing=1.15, space_after=4)
    card(s, Inches(6.83), Inches(2.3), Inches(5.9), Inches(1.55), fill=PANEL, line=IT_C)
    line_shape(s, Inches(6.83), Inches(2.3), Inches(0.09), Inches(1.55), IT_C)
    text(s, Inches(7.08), Inches(2.5), Inches(5.5), Inches(1.2),
         [[("FOR IT", SANS_SEMI, 11, IT_C, True, False, 1.5)],
          [("Architecture, data, security, integrations, and a step-by-step "
            "runbook to redeploy on Teroxx-owned cloud.", SANS, 12, INK_70)]],
         line_spacing=1.15, space_after=4)
    # contents list
    text(s, Inches(0.6), Inches(4.15), Inches(6), Inches(0.3),
         [[("WHAT'S INSIDE", SANS_SEMI, 10.5, ELECTRIC, True, False, 2.0)]])
    toc = [
        ("1", "What the app does", "Capabilities, workflow, the model & data — plain language"),
        ("2", "Project status", "What's live, known limitations, where it runs today"),
        ("3", "Architecture", "Stack, data, security, the proposal engine, integrations"),
        ("4", "Standing it up on your cloud", "Migration runbook, recommended path, credentials transfer"),
        ("5", "Ownership & handover", "Who owns what, first-week checklist, contacts & index"),
    ]
    y = 4.55
    for num, title, sub in toc:
        rect(s, Inches(0.6), Inches(y), Inches(0.42), Inches(0.42), fill=DEEP_INDIGO)
        text(s, Inches(0.6), Inches(y), Inches(0.42), Inches(0.42),
             [[(num, SANS_BOLD, 15, WHITE, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        text(s, Inches(1.2), Inches(y - 0.03), Inches(11.4), Inches(0.5),
             [[(title + "   ", SANS_SEMI, 13.5, INK, True), (sub, SANS, 11.5, INK_70)]],
             anchor=MSO_ANCHOR.MIDDLE)
        y += 0.5


# ═══════════════════════════════════════════════════════════════════════════
# SECTION DIVIDER
# ═══════════════════════════════════════════════════════════════════════════
def s_divider(num, title, items, chip=None):
    s = slide(); bg(s, NIGHTBLUE)
    rect(s, 0, 0, Inches(4.4), EMU_H, fill=DEEP_INDIGO)
    text(s, Inches(0.7), Inches(2.2), Inches(3), Inches(2),
         [[(num, SERIF, 120, EMBER)]])
    line_shape(s, Inches(4.9), Inches(2.5), Pt(2.5), Inches(2.4), EMBER)
    text(s, Inches(5.25), Inches(2.55), Inches(7.5), Inches(1.2),
         [[(title, SERIF, 40, WHITE)]], line_spacing=1.0)
    bl = []
    for it in items:
        bl.append([("—  ", SANS_SEMI, 14, EMBER, True), (it, SANS, 14, CREAM)])
    text(s, Inches(5.28), Inches(4.05), Inches(7.3), Inches(2.2), bl,
         line_spacing=1.1, space_after=7)
    if os.path.exists(LOGO_WHITE):
        s.shapes.add_picture(LOGO_WHITE, Inches(11.7), Inches(6.75), height=Inches(0.34))
    if chip:
        # light chip on dark
        label = {"research": "FOR RESEARCH", "it": "FOR IT", "both": "RESEARCH + IT"}[chip]
        col = {"research": SANDSTONE, "it": EMBER, "both": SANDSTONE}[chip]
        w = Inches(1.9)
        rect(s, Inches(5.28), Inches(1.95), w, Inches(0.36), fill=None, line=col, line_w=Pt(1.25))
        text(s, Inches(5.28), Inches(1.95), w, Inches(0.36),
             [[(label, SANS_SEMI, 9, col, True, False, 1.5)]],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE — EXECUTIVE SUMMARY
# ═══════════════════════════════════════════════════════════════════════════
def s_exec():
    s = slide(); content_header(s, "Section 1 · What the app does", "In one minute", chip="both")
    text(s, Inches(0.6), Inches(1.5), Inches(12.1), Inches(1.0),
         [[("The Portfolio Allocation App is the web version of Teroxx's Portfolio "
            "Allocation Model (v5.0). ", SANS_SEMI, 14.5, INK, True),
           ("An advisor picks a client's risk profile and investable universe; the app "
            "scores 79 crypto assets, builds a personalised allocation, tracks the live "
            "portfolio, and exports a fully branded client proposal as Word, PDF or Google Doc.",
            SANS, 14.5, INK_70)]], line_spacing=1.2)
    # metric strip
    metrics = [
        ("79", "assets scored", "multi-factor model"),
        ("8", "advisor tabs", "allocate → propose"),
        ("3", "export formats", "DOCX · PDF · GDoc"),
        ("Live", "in daily use", "by the advisory desk"),
    ]
    x = 0.6; w = 2.95; gap = 0.13
    for val, lab, sub in metrics:
        card(s, Inches(x), Inches(2.85), Inches(w), Inches(1.5), fill=PANEL, line=GRID)
        line_shape(s, Inches(x), Inches(2.85), Inches(w), Pt(3), ELECTRIC)
        text(s, Inches(x), Inches(3.02), Inches(w), Inches(0.7),
             [[(val, SERIF, 34, DEEP_INDIGO)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        text(s, Inches(x), Inches(3.72), Inches(w), Inches(0.6),
             [[(lab, SANS_SEMI, 11.5, INK, True)], [(sub, SANS, 10, INK_40)]],
             align=PP_ALIGN.CENTER, space_after=1)
        x += w + gap
    # the ask box
    card(s, Inches(0.6), Inches(4.7), Inches(12.13), Inches(1.85), fill=RGBColor(0xFB,0xF0,0xEC), line=EMBER)
    line_shape(s, Inches(0.6), Inches(4.7), Inches(0.09), Inches(1.85), EMBER)
    text(s, Inches(0.9), Inches(4.9), Inches(11.6), Inches(0.4),
         [[("WHY THIS HANDOVER MATTERS", SANS_SEMI, 11, EMBER, True, False, 1.8)]])
    text(s, Inches(0.9), Inches(5.28), Inches(11.6), Inches(1.2),
         [[("Today the app runs on Aleksander's ", SANS, 13, INK_70),
           ("private personal Render account", SANS_SEMI, 13, INK, True),
           (". When that account goes away, so does the deployment. The goal of this "
            "handover is to move it onto Teroxx-owned cloud and Teroxx-owned "
            "credentials, with Research owning the content and IT owning the infrastructure.",
            SANS, 13, INK_70)]], line_spacing=1.2)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE — THE 8 TABS
# ═══════════════════════════════════════════════════════════════════════════
def s_tabs():
    s = slide(); content_header(s, "Section 1 · What the app does", "The eight advisor tabs", chip="research")
    tabs = [
        ("Allocator", "Risk profile + universe → target allocation table & donut chart"),
        ("Your Portfolio", "Enter a portfolio value → position sizes with BUY actions"),
        ("Recurring Buys", "DCA planner: scope, horizon, minimum order + backtest"),
        ("Factor Scores", "The 5-factor model: weights, radar & composite scores"),
        ("Fundamentals", "The 10-factor fundamental model per asset"),
        ("Allocations", "Risk tilts, strategic weights, per-asset breakdown"),
        ("Rebalancing", "Current vs target → concrete BUY / SELL actions"),
        ("P&L Tracker", "Entry price/qty → unrealised P&L against live prices"),
    ]
    x0, y0, w, h, gx, gy = 0.6, 1.55, 5.97, 1.24, 0.19, 0.16
    for i, (t, d) in enumerate(tabs):
        r, c = divmod(i, 2)
        x = x0 + c * (w + gx); y = y0 + r * (h + gy)
        card(s, Inches(x), Inches(y), Inches(w), Inches(h), fill=PANEL, line=GRID)
        rect(s, Inches(x), Inches(y), Inches(0.5), Inches(h), fill=DEEP_INDIGO)
        text(s, Inches(x), Inches(y), Inches(0.5), Inches(h),
             [[(str(i + 1), SANS_BOLD, 17, WHITE, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        text(s, Inches(x + 0.66), Inches(y + 0.17), Inches(w - 0.8), Inches(1.0),
             [[(t, SANS_SEMI, 14, DEEP_INDIGO, True)], [(d, SANS, 11, INK_70)]],
             line_spacing=1.08, space_after=3)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE — CLIENT WORKFLOW
# ═══════════════════════════════════════════════════════════════════════════
def s_workflow():
    s = slide(); content_header(s, "Section 1 · What the app does", "From risk profile to client proposal", chip="research")
    text(s, Inches(0.6), Inches(1.45), Inches(12), Inches(0.5),
         [[("The everyday advisor journey — five steps, ending in a branded document ready to send.",
            SANS, 13, INK_70)]])
    steps = [
        ("1", "Select", "Choose the client, risk profile and investable universe"),
        ("2", "Allocate", "The app scores 79 assets and builds the target allocation"),
        ("3", "Review", "Check positions, rebalancing actions, P&L and drift"),
        ("4", "Personalise", "Add advisor notes, market view, execution plan (EN or DE)"),
        ("5", "Export", "One click → branded Word, PDF or Google Doc proposal"),
    ]
    x0, y, w, gap = 0.6, 2.35, 2.28, 0.15
    for i, (n, t, d) in enumerate(steps):
        x = x0 + i * (w + gap)
        card(s, Inches(x), Inches(y), Inches(w), Inches(2.6), fill=PANEL, line=GRID)
        line_shape(s, Inches(x), Inches(y), Inches(w), Pt(3), EMBER if i == 4 else ELECTRIC)
        rect(s, Inches(x + w/2 - 0.28), Inches(y + 0.3), Inches(0.56), Inches(0.56),
             fill=DEEP_INDIGO if i < 4 else EMBER)
        text(s, Inches(x + w/2 - 0.28), Inches(y + 0.3), Inches(0.56), Inches(0.56),
             [[(n, SANS_BOLD, 19, WHITE, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        text(s, Inches(x + 0.12), Inches(y + 1.02), Inches(w - 0.24), Inches(0.4),
             [[(t, SANS_SEMI, 14.5, DEEP_INDIGO, True)]], align=PP_ALIGN.CENTER)
        text(s, Inches(x + 0.15), Inches(y + 1.45), Inches(w - 0.3), Inches(1.1),
             [[(d, SANS, 11, INK_70)]], align=PP_ALIGN.CENTER, line_spacing=1.1)
        if i < 4:
            text(s, Inches(x + w - 0.02), Inches(y + 0.9), Inches(0.3), Inches(0.5),
                 [[("›", SANS_BOLD, 22, SANDSTONE, True)]], align=PP_ALIGN.CENTER)
    card(s, Inches(0.6), Inches(5.35), Inches(12.13), Inches(1.15), fill=RGBColor(0xF0,0xF5,0xF7), line=ELECTRIC)
    text(s, Inches(0.9), Inches(5.52), Inches(11.6), Inches(0.9),
         [[("Faithful proposals:  ", SANS_SEMI, 12.5, ELECTRIC, True),
           ("when the advisor hits export, the app freezes the exact on-screen allocation, so the "
            "document is a true record of what was reviewed — not a later recompute against moved "
            "market data. The Word file is the master; PDF and Google Doc are conversions of it, so the "
            "three can never drift apart.", SANS, 12.5, INK_70)]], line_spacing=1.15)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE — SELECTABLE UNIVERSES & BASKETS
# ═══════════════════════════════════════════════════════════════════════════
def s_universes():
    s = slide(); content_header(s, "Section 1 · What the app does", "The selectable universes & baskets", chip="research")
    text(s, Inches(0.6), Inches(1.42), Inches(12.1), Inches(0.5),
         [[("Before allocating, the advisor picks a starting list. Two kinds: the firm's ", SANS, 12.5, INK_70),
           ("model portfolios", SANS_SEMI, 12.5, INK, True),
           (" (broad, risk-tiered) or a single-theme ", SANS, 12.5, INK_70),
           ("thematic basket", SANS_SEMI, 12.5, INK, True),
           (" (index-style, 100% in-theme).", SANS, 12.5, INK_70)]], line_spacing=1.12)
    # LEFT — model portfolios
    text(s, Inches(0.6), Inches(2.05), Inches(6.5), Inches(0.3),
         [[("MODEL PORTFOLIOS  —  pick one", SANS_SEMI, 10.5, ELECTRIC, True, False, 1.5)]])
    models = [
        ("Teroxx Core", "9", "Live, active offering", DEEP_INDIGO,
         "USDC · EURC · PAXG · BTC · ETH · BNB · XRP · ADA · POL"),
        ("Teroxx Expanded", "21", "Confirmed MVP next step", ELECTRIC,
         "Core + MNT, and DeFi/infra: AAVE · UNI · COMP · EUL · PENDLE · SYRUP · ENA · ONDO · LINK · QNT · CHZ"),
        ("Teroxx Extended", "40", "Full acquirable universe", RGBColor(0x5C,0x6F,0xA1),
         "Expanded + majors & research picks: SOL · HYPE · TRX · AVAX · DOT · ARB · SUI · RENDER · AKT · AR · VVV · AERO · LTC · BCH · DOGE …"),
    ]
    y = 2.4
    for name, cnt, char, col, names in models:
        card(s, Inches(0.6), Inches(y), Inches(6.5), Inches(1.28), fill=PANEL, line=GRID)
        line_shape(s, Inches(0.6), Inches(y), Inches(0.09), Inches(1.28), col)
        rect(s, Inches(6.42), Inches(y + 0.16), Inches(0.55), Inches(0.5), fill=col)
        text(s, Inches(6.42), Inches(y + 0.16), Inches(0.55), Inches(0.5),
             [[(cnt, SANS_BOLD, 18, WHITE, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        text(s, Inches(0.82), Inches(y + 0.13), Inches(5.5), Inches(0.3),
             [[(name + "   ", SANS_SEMI, 13.5, DEEP_INDIGO, True), (char, SANS, 10.5, INK_40, False, True)]])
        text(s, Inches(0.82), Inches(y + 0.5), Inches(5.5), Inches(0.72),
             [[(names, SANS, 9.8, INK_70)]], line_spacing=1.08)
        y += 1.38
    # RIGHT — thematic baskets
    text(s, Inches(7.35), Inches(2.05), Inches(5.4), Inches(0.3),
         [[("THEMATIC BASKETS  —  index-style, in-theme", SANS_SEMI, 10.5, EMBER, True, False, 1.2)]])
    baskets = [
        ("DeFi", "11", "Lending, DEXs, perps & yield"),
        ("Smart-Contract L1", "9", "L1 platforms & settlement networks"),
        ("AI & Compute", "4", "Decentralised AI, GPU compute, storage"),
        ("RWA", "2", "Tokenised real-world assets & treasuries"),
        ("Layer 2", "3", "Ethereum scaling networks & rollups"),
        ("Payments", "2", "Established P2P payment networks"),
    ]
    y = 2.4
    for name, cnt, blurb in baskets:
        card(s, Inches(7.35), Inches(y), Inches(5.38), Inches(0.62), fill=PANEL, line=GRID)
        text(s, Inches(7.5), Inches(y), Inches(0.5), Inches(0.62),
             [[(cnt, SANS_BOLD, 15, EMBER, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        text(s, Inches(8.1), Inches(y), Inches(4.5), Inches(0.62),
             [[(name + "  Basket", SANS_SEMI, 11.5, DEEP_INDIGO, True)], [(blurb, SANS, 9.5, INK_70)]],
             anchor=MSO_ANCHOR.MIDDLE, space_after=0, line_spacing=1.0)
        y += 0.7
    # bottom strip
    card(s, Inches(0.6), Inches(6.5), Inches(12.13), Inches(0.02), fill=GRID, line=GRID)
    text(s, Inches(0.6), Inches(6.58), Inches(12.2), Inches(0.5),
         [[("Then choose:  ", SANS_SEMI, 10.5, ELECTRIC, True),
           ("risk profile ", SANS_SEMI, 10.5, INK, True),
           ("(Conservative · Balanced · Growth · Aggressive)  ×  ", SANS, 10.5, INK_70),
           ("allocation mode ", SANS_SEMI, 10.5, INK, True),
           ("(Standard = market-cap weighted · ", SANS, 10.5, INK_70),
           ("Fundamental = factor-score weighted, the default", SANS_SEMI, 10.5, INK, True),
           (").   The full research model scores 79 tokens; the 40-token union is what the firm can acquire today.",
            SANS, 10.5, INK_70)]],
         line_spacing=1.1)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE — THE MODEL
# ═══════════════════════════════════════════════════════════════════════════
def s_model():
    s = slide(); content_header(s, "Section 1 · What the app does", "How the allocation is built", chip="research")
    text(s, Inches(0.6), Inches(1.45), Inches(12), Inches(0.5),
         [[("Every asset earns a score; scores drive tier weights; a macro overlay tilts risk "
            "up or down. All transparent and rules-based.", SANS, 13, INK_70)]], line_spacing=1.15)
    cols = [
        ("SCORE EACH ASSET", ELECTRIC, [
            "Value-accrual signals: valuation (FDV/fees, FDV/TVL), fee revenue, dilution, buybacks",
            "Sector-tailored signals incl. price momentum & developer activity",
            "Every signal computed live from market & on-chain data",
        ]),
        ("BUILD THE TIERS", DEEP_INDIGO, [
            "Assets sorted into strategic tiers (core → satellite)",
            "Risk profile sets how much sits in each tier",
            "Defensive sleeve (stablecoins + gold hedge) for ballast",
        ]),
        ("OVERLAY THE REGIME", EMBER, [
            "22-indicator macro composite scores the market 0–100",
            "Sentiment, TradFi, on-chain and technical signals",
            "Tilts overall risk toward the current regime",
        ]),
    ]
    x0, w, gap = 0.6, 3.91, 0.2
    for i, (title, col, items) in enumerate(cols):
        x = x0 + i * (w + gap)
        card(s, Inches(x), Inches(2.2), Inches(w), Inches(3.5), fill=PANEL, line=GRID)
        rect(s, Inches(x), Inches(2.2), Inches(w), Inches(0.62), fill=col)
        text(s, Inches(x + 0.25), Inches(2.2), Inches(w - 0.4), Inches(0.62),
             [[(f"STEP {i+1}", SANS_SEMI, 9, WHITE, True, False, 1.5)],
              [(title, SANS_SEMI, 12.5, WHITE, True)]], anchor=MSO_ANCHOR.MIDDLE, space_after=0, line_spacing=1.0)
        bullets(s, Inches(x + 0.25), Inches(3.05), Inches(w - 0.5), Inches(2.5), items,
                size=11.5, gap=9, marker_color=col, lh=1.12)
    text(s, Inches(0.6), Inches(5.95), Inches(12), Inches(0.5),
         [[("Result: ", SANS_SEMI, 12, INK, True),
           ("a client-ready, risk-appropriate target allocation the advisor can explain line by line.",
            SANS, 12, INK_70)]])


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE — DATA SOURCES
# ═══════════════════════════════════════════════════════════════════════════
def s_datasources():
    s = slide(); content_header(s, "Section 1 · What the app does", "Where the numbers come from", chip="both")
    text(s, Inches(0.6), Inches(1.45), Inches(12), Inches(0.5),
         [[("The app pulls live market data from six external sources. Only two need an API key; "
            "the other four are public. Prices are cached to stay fast and within rate limits.", SANS, 13, INK_70)]], line_spacing=1.15)
    srcs = [
        ("CoinMarketCap", "Spot prices, market cap, 24h change", "API key", "5 min cache"),
        ("CoinGecko", "Fallback prices, dev activity, supply", "Demo key", "up to 24h"),
        ("Binance Futures", "Funding rates, open interest", "Public", "5 min cache"),
        ("DeFiLlama", "Protocol TVL, fees, DEX volume", "Public", "2 h cache"),
        ("Messari", "Network transactions & active addresses", "Public", "1 h cache"),
        ("Macro composite", "Fear/greed, VIX, treasuries, BTC trend", "Public mix", "hysteresis"),
    ]
    x0, y0, w, h, gx, gy = 0.6, 2.2, 3.91, 1.55, 0.2, 0.2
    for i, (name, use, auth, cache) in enumerate(srcs):
        r, c = divmod(i, 3)
        x = x0 + c * (w + gx); y = y0 + r * (h + gy)
        card(s, Inches(x), Inches(y), Inches(w), Inches(h), fill=PANEL, line=GRID)
        line_shape(s, Inches(x), Inches(y), Inches(0.08), Inches(h), ELECTRIC)
        text(s, Inches(x + 0.28), Inches(y + 0.16), Inches(w - 0.4), Inches(0.35),
             [[(name, SANS_SEMI, 13.5, DEEP_INDIGO, True)]])
        text(s, Inches(x + 0.28), Inches(y + 0.55), Inches(w - 0.45), Inches(0.6),
             [[(use, SANS, 11, INK_70)]], line_spacing=1.05)
        text(s, Inches(x + 0.28), Inches(y + 1.16), Inches(w - 0.45), Inches(0.3),
             [[("Auth: ", SANS_SEMI, 9.5, INK, True), (auth + "   ·   ", SANS, 9.5, INK_70),
               (cache, SANS, 9.5, ELECTRIC, True)]])
    text(s, Inches(0.6), Inches(6.18), Inches(12.2), Inches(0.7),
         [[("Action for IT:  ", SANS_SEMI, 11.5, EMBER, True),
           ("CoinMarketCap and CoinGecko must run on ", SANS, 11.5, INK_70),
           ("Teroxx-owned API keys", SANS_SEMI, 11.5, INK, True),
           (" — the app currently relies on shared / personal defaults. The four public sources need no key.",
            SANS, 11.5, INK_70)],
          [("If a source is briefly down, the app retries with backoff and falls back to cached values — it "
            "degrades gracefully rather than breaking.", SANS, 10.5, INK_40, False, True)]],
         line_spacing=1.1, space_after=3)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE — STATUS DASHBOARD
# ═══════════════════════════════════════════════════════════════════════════
def s_status():
    s = slide(); content_header(s, "Section 2 · Project status", "What's done, and what to watch", chip="both")
    # DONE column
    card(s, Inches(0.6), Inches(1.5), Inches(6.0), Inches(5.0), fill=RGBColor(0xEF,0xF6,0xF0), line=SUCCESS)
    text(s, Inches(0.85), Inches(1.7), Inches(5.5), Inches(0.4),
         [[("LIVE & WORKING", SANS_SEMI, 12, SUCCESS, True, False, 1.5)]])
    bullets(s, Inches(0.85), Inches(2.2), Inches(5.5), Inches(4.2), [
        ("Allocation engine — ", "all 8 tabs, 79 assets, both scoring models"),
        ("Client CRM — ", "clients, positions/lots, drift, scenarios, audit log"),
        ("Proposal export — ", "DOCX / PDF / Google Doc, EN + DE, fully branded"),
        ("Live market data — ", "6 sources with caching and graceful fallback"),
        ("Auth & audit — ", "session login; every change logged by user"),
        ("Stable v1 API — ", "token-authenticated client endpoints"),
        ("Deployed — ", "auto-deploys on every git push to main"),
    ], size=11.5, gap=8, marker="✓", marker_color=SUCCESS, lh=1.1)
    # WATCH column
    card(s, Inches(6.83), Inches(1.5), Inches(5.9), Inches(5.0), fill=RGBColor(0xFB,0xF0,0xEC), line=EMBER)
    text(s, Inches(7.08), Inches(1.7), Inches(5.4), Inches(0.4),
         [[("KNOWN LIMITATIONS & GOTCHAS", SANS_SEMI, 12, EMBER, True, False, 1.5)]])
    bullets(s, Inches(7.08), Inches(2.2), Inches(5.4), Inches(4.2), [
        ("Confirm the DB backend — ", "the app supports managed Postgres via DATABASE_URL; verify in the Render dashboard, and put the Teroxx deployment on managed Postgres (Section 4)."),
        ("No automated test suite — ", "changes are verified manually in the UI."),
        ("Free-tier memory — ", "vector charts are off to avoid out-of-memory on 512 MB; PNG charts used instead."),
        ("Google Docs geo-blocks — ", "export must run from an allowed region; some servers get 403."),
        ("Secrets in repo config — ", "a few keys sit in render.yaml; rotate them on migration."),
    ], size=11.5, gap=9, marker="!", marker_color=EMBER, lh=1.1)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE — WHERE IT LIVES TODAY
# ═══════════════════════════════════════════════════════════════════════════
def s_wherelives():
    s = slide(); content_header(s, "Section 2 · Project status", "Where it lives today", chip="it")
    text(s, Inches(0.6), Inches(1.45), Inches(12), Inches(0.5),
         [[("One running deployment, two code repositories. The deployment is the single point "
            "of failure this handover resolves.", SANS, 13, INK_70)]], line_spacing=1.15)
    # deployment card (risk)
    card(s, Inches(0.6), Inches(2.2), Inches(5.9), Inches(2.05), fill=RGBColor(0xFB,0xF0,0xEC), line=EMBER)
    line_shape(s, Inches(0.6), Inches(2.2), Inches(5.9), Pt(3), EMBER)
    text(s, Inches(0.85), Inches(2.4), Inches(5.4), Inches(0.35),
         [[("RUNTIME  —  action required", SANS_SEMI, 11, EMBER, True, False, 1.2)]])
    text(s, Inches(0.85), Inches(2.78), Inches(5.4), Inches(1.4),
         [[("Render.com", SANS_SEMI, 15, INK, True),
           ("  ·  free web-service tier", SANS, 12, INK_70)],
          [("Hosted on Aleksander's ", SANS, 12, INK_70),
           ("private personal account", SANS_SEMI, 12, DANGER, True),
           (". Docker image built from the repo; auto-deploys on push. This must move to a "
            "Teroxx-owned account/cloud.", SANS, 12, INK_70)]],
         line_spacing=1.18, space_after=5)
    # repos card
    card(s, Inches(6.83), Inches(2.2), Inches(5.9), Inches(2.05), fill=PANEL, line=GRID)
    line_shape(s, Inches(6.83), Inches(2.2), Inches(5.9), Pt(3), ELECTRIC)
    text(s, Inches(7.08), Inches(2.4), Inches(5.4), Inches(0.35),
         [[("SOURCE CODE  —  in good shape", SANS_SEMI, 11, ELECTRIC, True, False, 1.2)]])
    text(s, Inches(7.08), Inches(2.78), Inches(5.4), Inches(1.4),
         [[("GitLab (Teroxx)", SANS_SEMI, 13, INK, True),
           ("  — IT-created group; Aleksander is maintainer", SANS, 10.5, INK_70)],
          [("GitHub (private)", SANS_SEMI, 13, INK, True),
           ("  — personal mirror, to be retired", SANS, 10.5, INK_70)],
          [("Full codebase is on GitLab (", SANS, 11, INK_70), ("main", SANS_SEMI, 11, INK, True),
           (" + import). IT owns deployment going forward.", SANS, 11, INK_70)]],
         line_spacing=1.15, space_after=4)
    # what must be preserved
    card(s, Inches(0.6), Inches(4.5), Inches(12.13), Inches(2.0), fill=RGBColor(0xF0,0xF5,0xF7), line=ELECTRIC)
    text(s, Inches(0.85), Inches(4.68), Inches(11.6), Inches(0.35),
         [[("TO KEEP THE APP RUNNING AFTER MIGRATION, FOUR THINGS MUST TRANSFER", SANS_SEMI, 11, ELECTRIC, True, False, 1.2)]])
    items = [
        ("1  Code", "Already on Teroxx GitLab — nothing to move"),
        ("2  A host", "Any Docker-capable cloud (Section 4)"),
        ("3  A database", "A Teroxx-managed Postgres (RDS); verify current backend"),
        ("4  Credentials", "API keys, session secret, Google service account"),
    ]
    x = 0.85; w = 2.9
    for t, d in items:
        text(s, Inches(x), Inches(5.15), Inches(w), Inches(1.2),
             [[(t, SANS_SEMI, 13, DEEP_INDIGO, True)], [(d, SANS, 11, INK_70)]],
             line_spacing=1.1, space_after=3)
        x += w + 0.13


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE — TECH STACK / ARCHITECTURE DIAGRAM
# ═══════════════════════════════════════════════════════════════════════════
def s_architecture():
    s = slide(); content_header(s, "Section 3 · Architecture", "How the system fits together", chip="it")
    # Three layers: Client -> App container -> External
    def box(x, y, w, h, title, sub, fill, tcol=WHITE, subcol=None):
        card(s, Inches(x), Inches(y), Inches(w), Inches(h), fill=fill, line=fill)
        text(s, Inches(x + 0.1), Inches(y + 0.1), Inches(w - 0.2), Inches(h - 0.2),
             [[(title, SANS_SEMI, 11.5, tcol, True)]] + ([[(sub, SANS, 9, subcol or tcol)]] if sub else []),
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=2, line_spacing=1.0)

    # Browser
    box(0.6, 1.7, 2.1, 1.0, "Advisor's browser", "HTMX + Chart.js UI", HEADER_BG)
    # arrow
    text(s, Inches(2.72), Inches(2.0), Inches(0.5), Inches(0.4), [[("→", SANS_BOLD, 20, INK_40, True)]], align=PP_ALIGN.CENTER)
    # App container (big)
    card(s, Inches(3.35), Inches(1.55), Inches(6.4), Inches(4.55), fill=PANEL, line=DEEP_INDIGO)
    text(s, Inches(3.35), Inches(1.65), Inches(6.4), Inches(0.35),
         [[("DOCKER CONTAINER  —  the whole app", SANS_SEMI, 10.5, DEEP_INDIGO, True, False, 1.2)]], align=PP_ALIGN.CENTER)
    box(3.6, 2.15, 5.9, 0.85, "FastAPI + Uvicorn", "web server · routes · session auth · v1 API", DEEP_INDIGO)
    box(3.6, 3.12, 2.85, 0.95, "Allocation engine", "scoring · DCA · rebalance · P&L", ELECTRIC)
    box(6.65, 3.12, 2.85, 0.95, "Proposal engine", "DOCX master → PDF / GDoc", ELECTRIC)
    box(3.6, 4.2, 2.85, 0.85, "Jinja templates", "server-rendered HTML", RGBColor(0x5C,0x6F,0xA1))
    box(6.65, 4.2, 2.85, 0.85, "LibreOffice + fonts", "DOCX→PDF · Söhne baked in", RGBColor(0x5C,0x6F,0xA1))
    box(3.6, 5.15, 5.9, 0.8, "In-memory caches", "prices · macro · DeFi (warm on startup)", SANDSTONE, tcol=INK)
    # External services (right)
    box(10.15, 1.7, 2.55, 0.95, "PostgreSQL", "clients · lots · audit · snapshots", INK)
    text(s, Inches(9.75), Inches(2.0), Inches(0.45), Inches(0.4), [[("↔", SANS_BOLD, 18, INK_40, True)]], align=PP_ALIGN.CENTER)
    ext = [
        ("Market data APIs", "CMC · CoinGecko · Binance"),
        ("On-chain data", "DeFiLlama · Messari"),
        ("Google Drive API", "Google Doc export"),
    ]
    y = 2.95
    for t, d in ext:
        box(10.15, y, 2.55, 0.95, t, d, HEADER_BG)
        text(s, Inches(9.75), Inches(y + 0.28), Inches(0.45), Inches(0.4), [[("←", SANS_BOLD, 16, INK_40, True)]], align=PP_ALIGN.CENTER)
        y += 1.07
    text(s, Inches(0.6), Inches(6.32), Inches(12.2), Inches(0.5),
         [[("No dependencies on other internal Teroxx systems", SANS_SEMI, 11.5, DEEP_INDIGO, True),
           (" (confirmed with IT) — only outbound calls to public data APIs. Everything else ships inside one "
            "container image, which is what makes it portable to any cloud.", SANS, 11.5, INK_70, False, True)]],
         line_spacing=1.1)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE — DATA & PERSISTENCE
# ═══════════════════════════════════════════════════════════════════════════
def s_data():
    s = slide(); content_header(s, "Section 3 · Architecture", "Data & persistence", chip="it")
    text(s, Inches(0.6), Inches(1.45), Inches(12), Inches(0.5),
         [[("SQLAlchemy ORM. The connection string is chosen at startup: ", SANS, 12.5, INK_70),
           ("DATABASE_URL", SANS_SEMI, 12.5, ELECTRIC, True),
           (" (Postgres) → ", SANS, 12.5, INK_70),
           ("TEROXX_DB_PATH", SANS_SEMI, 12.5, ELECTRIC, True),
           (" → local SQLite → /tmp fallback.", SANS, 12.5, INK_70)]], line_spacing=1.15)
    # models
    card(s, Inches(0.6), Inches(2.15), Inches(6.0), Inches(4.35), fill=PANEL, line=GRID)
    text(s, Inches(0.85), Inches(2.32), Inches(5.5), Inches(0.35),
         [[("SIX TABLES", SANS_SEMI, 11, ELECTRIC, True, False, 1.5)]])
    bullets(s, Inches(0.85), Inches(2.78), Inches(5.5), Inches(3.6), [
        ("Client — ", "profile, domicile, currency, defaults, overrides"),
        ("ClientLot — ", "individual holdings: ticker, qty, entry price/date"),
        ("AllocationSnapshot — ", "frozen allocations behind each proposal"),
        ("AdvisorAction — ", "full audit trail: who did what, when"),
        ("ApiToken — ", "hashed tokens for the v1 API"),
        ("Scenario — ", "saved A/B allocation comparisons"),
    ], size=11.5, gap=10, marker="▪", marker_color=DEEP_INDIGO, lh=1.1)
    # persistence note
    card(s, Inches(6.83), Inches(2.15), Inches(5.9), Inches(4.35), fill=RGBColor(0xFB,0xF0,0xEC), line=EMBER)
    text(s, Inches(7.08), Inches(2.32), Inches(5.4), Inches(0.35),
         [[("CONFIRM THE BACKEND", SANS_SEMI, 11, ELECTRIC, True, False, 1.5)]])
    text(s, Inches(7.08), Inches(2.78), Inches(5.4), Inches(3.5),
         [[("The app reads ", SANS, 12, INK_70), ("DATABASE_URL", SANS_SEMI, 12, INK, True),
           (" at startup. If it points at a managed Postgres (likely set in the Render dashboard), "
            "data is durable. If it is unset, the app falls back to SQLite on Render's free tier, whose "
            "disk is ", SANS, 12, INK_70), ("ephemeral", SANS_SEMI, 12, EMBER, True),
           (".", SANS, 12, INK_70)],
          [("", SANS, 6, INK_70)],
          [("Action: ", SANS_SEMI, 12, INK, True),
           ("confirm the current backend in the Render dashboard. Client data is mostly test/throwaway, "
            "so this is not blocking — a fresh database is fine.", SANS, 12, INK_70)],
          [("", SANS, 6, INK_70)],
          [("For Teroxx: ", SANS_SEMI, 12, SUCCESS, True),
           ("deploy on a managed Postgres (RDS) with DATABASE_URL set, and enable backups. The code "
            "already supports it — no changes needed (step 2 of the runbook).", SANS, 12, INK_70)]],
         line_spacing=1.16, space_after=2)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE — AUTH & SECURITY
# ═══════════════════════════════════════════════════════════════════════════
def s_security():
    s = slide(); content_header(s, "Section 3 · Architecture", "Access & security model", chip="it")
    cols = [
        ("WHO CAN LOG IN", ELECTRIC, [
            "Session-based login, no external identity provider",
            "Four users today: Jannick, Leonardo, Aleksander, Christopher",
            "Passwords stored salted + SHA-256 hashed in auth.py",
            "Emil (IT) has no login yet — add on handover",
        ]),
        ("HOW SESSIONS ARE SIGNED", DEEP_INDIGO, [
            "Signed cookies via SESSION_SECRET (7-day expiry)",
            "Default dev secret MUST be overridden in production",
            "Middleware guards all UI routes; /health & static are open",
            "Regenerate the secret on migration (logs everyone out once)",
        ]),
        ("MACHINE ACCESS (v1 API)", EMBER, [
            "Bearer tokens for /api/v1/* integrations",
            "Tokens stored hashed; plaintext shown once at creation",
            "Scoped (read / write / admin); revocable",
            "Lost tokens can't be recovered — mint a new one",
        ]),
    ]
    x0, w, gap = 0.6, 3.91, 0.2
    for i, (t, col, items) in enumerate(cols):
        x = x0 + i * (w + gap)
        card(s, Inches(x), Inches(1.55), Inches(w), Inches(4.35), fill=PANEL, line=GRID)
        rect(s, Inches(x), Inches(1.55), Inches(w), Inches(0.58), fill=col)
        text(s, Inches(x + 0.22), Inches(1.55), Inches(w - 0.4), Inches(0.58),
             [[(t, SANS_SEMI, 11.5, WHITE, True)]], anchor=MSO_ANCHOR.MIDDLE)
        bullets(s, Inches(x + 0.22), Inches(2.35), Inches(w - 0.44), Inches(3.4), items,
                size=11, gap=10, marker_color=col, lh=1.12)
    text(s, Inches(0.6), Inches(6.1), Inches(12.13), Inches(0.6),
         [[("Note for IT:  ", SANS_SEMI, 11.5, EMBER, True),
           ("this is lightweight internal auth, adequate for a small trusted desk. If the app is ever "
            "exposed more widely, front it with SSO / an identity-aware proxy and move users out of source code.",
            SANS, 11.5, INK_70)]], line_spacing=1.12)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE — PROPOSAL ENGINE
# ═══════════════════════════════════════════════════════════════════════════
def s_proposal_engine():
    s = slide(); content_header(s, "Section 3 · Architecture", "The proposal engine (and its dependencies)", chip="it")
    text(s, Inches(0.6), Inches(1.45), Inches(12), Inches(0.5),
         [[("One design rule keeps the three outputs identical: ", SANS, 12.5, INK_70),
           ("the Word document is the single source of truth.", SANS_SEMI, 12.5, INK, True)]], line_spacing=1.15)
    # pipeline
    stages = [
        ("Build context", "proposal.py assembles allocation + market data + narrative (EN/DE)", ELECTRIC),
        ("Render DOCX", "proposal_docx.py — python-docx, brand fonts, charts as images", DEEP_INDIGO),
        ("Convert PDF", "LibreOffice (soffice) renders the DOCX to PDF", RGBColor(0x5C,0x6F,0xA1)),
        ("Or Google Doc", "upload DOCX to Drive → native Google Doc", EMBER),
    ]
    x0, y, w, gap = 0.6, 2.25, 2.95, 0.13
    for i, (t, d, col) in enumerate(stages):
        x = x0 + i * (w + gap)
        card(s, Inches(x), Inches(y), Inches(w), Inches(1.75), fill=PANEL, line=GRID)
        line_shape(s, Inches(x), Inches(y), Inches(w), Pt(3), col)
        text(s, Inches(x + 0.2), Inches(y + 0.18), Inches(w - 0.4), Inches(0.4),
             [[(t, SANS_SEMI, 13.5, DEEP_INDIGO, True)]])
        text(s, Inches(x + 0.2), Inches(y + 0.62), Inches(w - 0.4), Inches(1.0),
             [[(d, SANS, 11, INK_70)]], line_spacing=1.12)
        if i < 3:
            text(s, Inches(x + w - 0.03), Inches(y + 0.55), Inches(0.25), Inches(0.5),
                 [[("›", SANS_BOLD, 20, SANDSTONE, True)]], align=PP_ALIGN.CENTER)
    # deps box
    card(s, Inches(0.6), Inches(4.3), Inches(12.13), Inches(2.2), fill=RGBColor(0xF0,0xF5,0xF7), line=ELECTRIC)
    text(s, Inches(0.85), Inches(4.48), Inches(11.6), Inches(0.35),
         [[("WHAT THE CONTAINER MUST PROVIDE (all baked into the Docker image already)", SANS_SEMI, 11, ELECTRIC, True, False, 1.0)]])
    deps = [
        ("LibreOffice Writer", "converts DOCX → PDF (the soffice binary)"),
        ("Pango / Cairo / HarfBuzz", "rasterise SVG charts and shape text"),
        ("Brand fonts", "Söhne + Sometimes Times installed system-wide"),
        ("TEROXX_VECTOR_CHARTS=0", "PNG charts on small instances to avoid OOM"),
    ]
    x = 0.85; w = 2.95
    for t, d in deps:
        text(s, Inches(x), Inches(4.95), Inches(w - 0.05), Inches(1.4),
             [[(t, SANS_SEMI, 12, DEEP_INDIGO, True)], [(d, SANS, 10.5, INK_70)]],
             line_spacing=1.1, space_after=3)
        x += w + 0.1


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE — ENV VAR REFERENCE TABLE
# ═══════════════════════════════════════════════════════════════════════════
def s_envvars():
    s = slide(); content_header(s, "Section 3 · Architecture", "Environment-variable reference", chip="it")
    text(s, Inches(0.6), Inches(1.42), Inches(12), Inches(0.4),
         [[("Everything configurable is an environment variable. This is the full contract between the app and its host.",
            SANS, 12, INK_70)]])
    rows = [
        ("DATABASE_URL", "Postgres connection string", "Set on migration", "●"),
        ("SESSION_SECRET", "Signs login cookies", "Regenerate — secret", "●"),
        ("CMC_API_KEY", "CoinMarketCap prices", "Has public default", "○"),
        ("COINGECKO_API_KEY", "CoinGecko fallback data", "Has demo default", "○"),
        ("GOOGLE_SERVICE_ACCOUNT_JSON", "Google Doc export credentials", "Feature-off if unset", "◐"),
        ("GOOGLE_DRIVE_FOLDER_ID", "Target Drive/Shared-Drive folder", "Optional", "◐"),
        ("GOOGLE_IMPERSONATE_SUBJECT", "Domain-wide delegation user", "Optional", "◐"),
        ("GOOGLE_DOCS_SHARE", "Share mode: none / domain / anyone", "Default none", "◐"),
        ("TEROXX_VECTOR_CHARTS", "Vector(1)/PNG(0) charts", "0 on small hosts", "○"),
        ("PORT", "Port the server binds to", "Host usually sets", "○"),
    ]
    # header
    tx, ty, tw = 0.6, 2.05, 12.13
    c1, c2, c3, c4 = 3.5, 4.9, 3.0, 0.73
    hdr = rect(s, Inches(tx), Inches(ty), Inches(tw), Inches(0.42), fill=DEEP_INDIGO)
    heads = [("VARIABLE", c1), ("PURPOSE", c2), ("ON MIGRATION", c3), ("REQ", c4)]
    xx = tx
    for h, cw in heads:
        text(s, Inches(xx + 0.12), Inches(ty), Inches(cw - 0.15), Inches(0.42),
             [[(h, SANS_SEMI, 9.5, WHITE, True, False, 1.0)]], anchor=MSO_ANCHOR.MIDDLE)
        xx += cw
    y = ty + 0.42
    rh = 0.402
    for i, (v, p, m, req) in enumerate(rows):
        if i % 2 == 0:
            rect(s, Inches(tx), Inches(y), Inches(tw), Inches(rh), fill=PANEL, line=PANEL)
        cells = [(v, c1, "code"), (p, c2, "body"), (m, c3, "body"), (req, c4, "req")]
        xx = tx
        for val, cw, kind in cells:
            if kind == "code":
                text(s, Inches(xx + 0.12), Inches(y), Inches(cw - 0.15), Inches(rh),
                     [[(val, SANS_SEMI, 10.5, DEEP_INDIGO, True)]], anchor=MSO_ANCHOR.MIDDLE)
            elif kind == "req":
                col = DANGER if val == "●" else (EMBER if val == "◐" else INK_40)
                lab = {"●": "required", "◐": "optional", "○": "default ok"}[val]
                text(s, Inches(xx + 0.12), Inches(y), Inches(cw - 0.1), Inches(rh),
                     [[(lab, SANS_SEMI, 8.5, col, True)]], anchor=MSO_ANCHOR.MIDDLE)
            else:
                text(s, Inches(xx + 0.12), Inches(y), Inches(cw - 0.15), Inches(rh),
                     [[(val, SANS, 10.5, INK_70)]], anchor=MSO_ANCHOR.MIDDLE)
            xx += cw
        y += rh
    text(s, Inches(0.6), Inches(y + 0.08), Inches(12), Inches(0.35),
         [[("● required for a healthy production deploy    ◐ needed only for Google Doc export    ○ safe default exists",
            SANS, 9.5, INK_40)]])


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE — MIGRATION OVERVIEW
# ═══════════════════════════════════════════════════════════════════════════
def s_migration_overview():
    s = slide(); content_header(s, "Section 4 · Standing it up on your cloud", "The migration in one picture", chip="it")
    text(s, Inches(0.6), Inches(1.45), Inches(12), Inches(0.7),
         [[("The app is a ", SANS, 13, INK_70), ("single Docker container", SANS_SEMI, 13, INK, True),
           (" plus a ", SANS, 13, INK_70), ("Postgres database", SANS_SEMI, 13, INK, True),
           (" plus a handful of ", SANS, 13, INK_70), ("environment variables", SANS_SEMI, 13, INK, True),
           (". Anything that can run a container can run this — no code changes required.", SANS, 13, INK_70)]],
         line_spacing=1.18)
    pieces = [
        ("The container", "Built from the repo's Dockerfile. Bundles the app, LibreOffice and brand fonts. Nothing else to install.", ELECTRIC),
        ("The database", "A managed PostgreSQL. Point the app at it with DATABASE_URL. Durable, backed up, no ops.", DEEP_INDIGO),
        ("The secrets", "~5 environment variables (Section 3 table). Set once in the host's dashboard.", EMBER),
        ("The domain", "Optional: map a Teroxx URL and TLS certificate. Most hosts do this in a click.", RGBColor(0x5C,0x6F,0xA1)),
    ]
    x0, y0, w, h, gx, gy = 0.6, 2.5, 5.97, 1.85, 0.19, 0.18
    for i, (t, d, col) in enumerate(pieces):
        r, c = divmod(i, 2)
        x = x0 + c * (w + gx); y = y0 + r * (h + gy)
        card(s, Inches(x), Inches(y), Inches(w), Inches(h), fill=PANEL, line=GRID)
        line_shape(s, Inches(x), Inches(y), Inches(0.09), Inches(h), col)
        text(s, Inches(x + 0.3), Inches(y + 0.2), Inches(w - 0.5), Inches(0.4),
             [[(t, SANS_SEMI, 15, DEEP_INDIGO, True)]])
        text(s, Inches(x + 0.3), Inches(y + 0.65), Inches(w - 0.55), Inches(1.1),
             [[(d, SANS, 12, INK_70)]], line_spacing=1.15)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE — RUNBOOK (step by step)
# ═══════════════════════════════════════════════════════════════════════════
def s_runbook():
    s = slide(); content_header(s, "Section 4 · Standing it up on your cloud", "Step-by-step deployment runbook", chip="it")
    steps = [
        ("1", "Get the code", "Clone from Teroxx GitLab. IT should own the repo and CI going forward."),
        ("2", "Provision Postgres", "Create a managed PostgreSQL; copy its connection string."),
        ("3", "Provision the host", "Create a container service pointing at the repo's Dockerfile."),
        ("4", "Set env vars", "DATABASE_URL, a fresh SESSION_SECRET, API keys, TEROXX_VECTOR_CHARTS."),
        ("5", "Deploy", "Build the image and start it. Tables auto-create on first boot."),
        ("6", "Verify health", "Hit /health (200 OK), log in, generate a test PDF proposal."),
        ("7", "Map the domain", "Point a Teroxx subdomain at the service; enable TLS."),
        ("8", "Cut over", "Update bookmarks, add Emil's login, suspend the old Render app."),
    ]
    x0, y0, w, h, gx, gy = 0.6, 1.5, 5.97, 1.12, 0.19, 0.13
    for i, (n, t, d) in enumerate(steps):
        r, c = divmod(i, 2)
        x = x0 + c * (w + gx); y = y0 + r * (h + gy)
        card(s, Inches(x), Inches(y), Inches(w), Inches(h), fill=PANEL, line=GRID)
        rect(s, Inches(x), Inches(y), Inches(0.5), Inches(h), fill=DEEP_INDIGO if i < 7 else EMBER)
        text(s, Inches(x), Inches(y), Inches(0.5), Inches(h),
             [[(n, SANS_BOLD, 16, WHITE, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        text(s, Inches(x + 0.64), Inches(y + 0.15), Inches(w - 0.78), Inches(h - 0.2),
             [[(t, SANS_SEMI, 12.5, DEEP_INDIGO, True)], [(d, SANS, 10, INK_70)]],
             line_spacing=1.05, space_after=2, anchor=MSO_ANCHOR.MIDDLE)
    text(s, Inches(0.6), Inches(6.62), Inches(12), Inches(0.35),
         [[("Because the whole runtime is in the image, steps 3–5 are the same on every provider — only the dashboard differs. "
            "A companion runbook with exact commands accompanies this deck.", SANS, 10.5, INK_40, False, True)]])


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE — RECOMMENDED PATH & OPTIONS
# ═══════════════════════════════════════════════════════════════════════════
def s_hosting_options():
    s = slide(); content_header(s, "Section 4 · Standing it up on your cloud", "Where to run it — a recommendation", chip="it")
    # recommended
    card(s, Inches(0.6), Inches(1.5), Inches(12.13), Inches(1.65), fill=RGBColor(0xEF,0xF6,0xF0), line=SUCCESS)
    line_shape(s, Inches(0.6), Inches(1.5), Inches(0.1), Inches(1.65), SUCCESS)
    text(s, Inches(0.9), Inches(1.68), Inches(11.5), Inches(0.35),
         [[("RECOMMENDED — AWS, aligned with Teroxx IT", SANS_SEMI, 12, SUCCESS, True, False, 1.2)]])
    text(s, Inches(0.9), Inches(2.08), Inches(11.6), Inches(1.0),
         [[("IT has indicated ", SANS, 12.5, INK_70),
           ("AWS with the EKS container platform", SANS_SEMI, 12.5, INK, True),
           (" as Teroxx's standard, plus a managed ", SANS, 12.5, INK_70),
           ("RDS Postgres", SANS_SEMI, 12.5, INK, True),
           (". The app is a single small, stateless service, so on EKS it is just one Deployment + Service + "
            "Ingress. If lower operational overhead is preferred, the same image runs unchanged on ", SANS, 12.5, INK_70),
           ("AWS App Runner or ECS Fargate", SANS_SEMI, 12.5, INK, True),
           (". Exact setup to be finalised with IT.", SANS, 12.5, INK_70)]], line_spacing=1.18)
    text(s, Inches(0.6), Inches(3.35), Inches(12), Inches(0.3),
         [[("OPTIONS AT A GLANCE", SANS_SEMI, 10.5, ELECTRIC, True, False, 1.8)]])
    opts = [
        ("AWS EKS\n(Teroxx container standard)", "One Deployment + Service + Ingress, with RDS Postgres alongside. IT's named platform.", "Target — align with IT", SUCCESS),
        ("AWS App Runner / ECS Fargate", "Same image, AWS-native, far less to operate than a full cluster.", "Lower-ops AWS option", ELECTRIC),
        ("PaaS\n(Render / Railway / Fly.io)", "Same model as today; render.yaml already in the repo. A quick bridge.", "Fastest lift-and-shift", EMBER),
        ("Other clouds\n(Cloud Run / Container Apps)", "The identical container runs anywhere — noted for completeness.", "Portable if ever needed", INK_70),
    ]
    x0, y, w, gap = 0.6, 3.75, 2.95, 0.13
    for t, d, tag, col in opts:
        x = x0
        card(s, Inches(x0), Inches(y), Inches(w), Inches(2.6), fill=PANEL, line=GRID)
        rect(s, Inches(x0), Inches(y), Inches(w), Pt(3.5), fill=col)
        text(s, Inches(x0 + 0.2), Inches(y + 0.16), Inches(w - 0.4), Inches(0.85),
             [[(t, SANS_SEMI, 11.5, DEEP_INDIGO, True)]], line_spacing=1.0)
        text(s, Inches(x0 + 0.2), Inches(y + 1.02), Inches(w - 0.4), Inches(1.05),
             [[(d, SANS, 10.5, INK_70)]], line_spacing=1.12)
        text(s, Inches(x0 + 0.2), Inches(y + 2.12), Inches(w - 0.4), Inches(0.4),
             [[(tag, SANS_SEMI, 10, col, True)]])
        x0 += w + gap


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE — CREDENTIALS & DATA TRANSFER
# ═══════════════════════════════════════════════════════════════════════════
def s_credentials():
    s = slide(); content_header(s, "Section 4 · Standing it up on your cloud", "Credentials & data transfer checklist", chip="it")
    text(s, Inches(0.6), Inches(1.42), Inches(12), Inches(0.45),
         [[("No live secrets are printed in this deck. ", SANS_SEMI, 12, INK, True),
           ("Aleksander will hand these over through a secure channel (password manager / vault), not email or chat.",
            SANS, 12, INK_70)]], line_spacing=1.12)
    # credentials
    card(s, Inches(0.6), Inches(2.05), Inches(6.0), Inches(4.45), fill=PANEL, line=GRID)
    text(s, Inches(0.85), Inches(2.22), Inches(5.5), Inches(0.35),
         [[("SECRETS TO TRANSFER & ROTATE", SANS_SEMI, 11, EMBER, True, False, 1.3)]])
    bullets(s, Inches(0.85), Inches(2.68), Inches(5.5), Inches(3.7), [
        ("SESSION_SECRET — ", "generate a brand-new random value on the new host"),
        ("CMC_API_KEY — ", "register Teroxx's own CoinMarketCap account & key; stop using the shared default"),
        ("COINGECKO_API_KEY — ", "register Teroxx's own CoinGecko key; stop using the personal one"),
        ("Google service account — ", "re-issue under Teroxx Workspace; share the Drive folder"),
        ("Repo access — ", "transfer GitLab ownership to IT; retire the personal GitHub mirror"),
        ("Render account — ", "hand over or rebuild; then delete the personal deployment"),
    ], size=11, gap=9, marker="→", marker_color=EMBER, lh=1.12)
    # data migration
    card(s, Inches(6.83), Inches(2.05), Inches(5.9), Inches(4.45), fill=RGBColor(0xF0,0xF5,0xF7), line=ELECTRIC)
    text(s, Inches(7.08), Inches(2.22), Inches(5.4), Inches(0.35),
         [[("MOVING EXISTING DATA", SANS_SEMI, 11, ELECTRIC, True, False, 1.3)]])
    text(s, Inches(7.08), Inches(2.68), Inches(5.4), Inches(3.6),
         [[("Client data is ", SANS, 12, INK_70),
           ("mostly test / throwaway", SANS_SEMI, 12, INK, True),
           (", so migration can start with a fresh Teroxx database — clients are re-entered as needed.",
            SANS, 12, INK_70)],
          [("", SANS, 5, INK_70)],
          [("If anything must be kept: ", SANS_SEMI, 12, INK, True),
           ("confirm the current backend, then dump it (Postgres or SQLite) and restore into the new "
            "Teroxx Postgres before cut-over. The schema is identical — SQLAlchemy creates it automatically.",
            SANS, 12, INK_70)],
          [("", SANS, 5, INK_70)],
          [("Going forward: ", SANS_SEMI, 12, SUCCESS, True),
           ("managed Postgres + automated daily backups keeps client data durable and recoverable.",
            SANS, 12, INK_70)]],
         line_spacing=1.16, space_after=2)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE — OWNERSHIP
# ═══════════════════════════════════════════════════════════════════════════
def s_ownership():
    s = slide(); content_header(s, "Section 5 · Ownership & handover", "Who owns what, going forward", chip="both")
    text(s, Inches(0.6), Inches(1.45), Inches(12), Inches(0.5),
         [[("A clean split: Research owns the content and daily use; IT owns the platform it runs on.",
            SANS, 13, INK_70)]])
    # two owner columns
    card(s, Inches(0.6), Inches(2.15), Inches(6.0), Inches(4.35), fill=PANEL, line=RESEARCH_C)
    rect(s, Inches(0.6), Inches(2.15), Inches(6.0), Inches(0.62), fill=RESEARCH_C)
    text(s, Inches(0.85), Inches(2.15), Inches(5.5), Inches(0.62),
         [[("RESEARCH  —  Jannick & Leonardo", SANS_SEMI, 13, WHITE, True)]], anchor=MSO_ANCHOR.MIDDLE)
    bullets(s, Inches(0.85), Inches(3.0), Inches(5.5), Inches(3.4), [
        "Day-to-day use: allocations, proposals, client reviews",
        "The model: universe, factor weights, tier rules, tilts",
        "Proposal content: narrative templates, EN/DE wording",
        "Deciding what changes are needed and prioritising them",
        "Adding / offboarding client records and advisor users",
    ], size=12, gap=11, marker_color=RESEARCH_C, lh=1.12)
    card(s, Inches(6.83), Inches(2.15), Inches(5.9), Inches(4.35), fill=PANEL, line=IT_C)
    rect(s, Inches(6.83), Inches(2.15), Inches(5.9), Inches(0.62), fill=IT_C)
    text(s, Inches(7.08), Inches(2.15), Inches(5.4), Inches(0.62),
         [[("IT  —  Emil", SANS_SEMI, 13, WHITE, True)]], anchor=MSO_ANCHOR.MIDDLE)
    bullets(s, Inches(7.08), Inches(3.0), Inches(5.4), Inches(3.4), [
        "The cloud host, the database and its backups",
        "Secrets & environment variables (rotation, storage)",
        "The repo, CI/CD and deployments on push",
        "Domain, TLS, uptime and monitoring /health",
        "Applying code changes Research requests (or a dev)",
    ], size=12, gap=11, marker_color=IT_C, lh=1.12)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE — FIRST-WEEK CHECKLIST + INDEX
# ═══════════════════════════════════════════════════════════════════════════
def s_checklist():
    s = slide(); content_header(s, "Section 5 · Ownership & handover", "First-week checklist & key references", chip="both")
    # checklist
    card(s, Inches(0.6), Inches(1.5), Inches(6.0), Inches(5.0), fill=RGBColor(0xF0,0xF5,0xF7), line=ELECTRIC)
    text(s, Inches(0.85), Inches(1.68), Inches(5.5), Inches(0.35),
         [[("FIRST WEEK — IN ORDER", SANS_SEMI, 11.5, ELECTRIC, True, False, 1.3)]])
    bullets(s, Inches(0.85), Inches(2.18), Inches(5.5), Inches(4.2), [
        "IT takes ownership of the GitLab repo",
        "Provision Postgres + container host (Section 4)",
        "Securely receive & set the secrets; rotate keys",
        "Deploy; verify /health, login, a test PDF proposal",
        "Map a Teroxx domain + TLS",
        "Add Emil's login; confirm Research can log in",
        "Cut over, then suspend the personal Render app",
        "Turn on database backups",
    ], size=11.5, gap=8.5, marker="☐", marker_color=DEEP_INDIGO, lh=1.1)
    # references / index
    card(s, Inches(6.83), Inches(1.5), Inches(5.9), Inches(5.0), fill=PANEL, line=GRID)
    text(s, Inches(7.08), Inches(1.68), Inches(5.4), Inches(0.35),
         [[("KEY FILES IN THE REPO", SANS_SEMI, 11.5, EMBER, True, False, 1.3)]])
    refs = [
        ("Dockerfile", "the entire runtime recipe"),
        ("render.yaml", "current deploy config & env keys"),
        ("requirements.txt", "Python dependencies"),
        ("app/db.py", "database models & connection logic"),
        ("app/auth.py", "users & session secret"),
        ("app/main.py", "all routes / the API surface"),
        ("app/pdf/", "the proposal (DOCX/PDF/GDoc) engine"),
        ("CLAUDE.md", "in-repo notes & design decisions"),
        ("handover/", "this deck + companion runbook"),
    ]
    tb = s.shapes.add_textbox(Inches(7.08), Inches(2.18), Inches(5.4), Inches(4.2))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left=0; tf.margin_right=0; tf.margin_top=0; tf.margin_bottom=0
    for i, (f, d) in enumerate(refs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = 1.1; p.space_after = Pt(8.5)
        r1 = p.add_run(); _set_run(r1, f, SANS_SEMI, 11.5, DEEP_INDIGO, True)
        r2 = p.add_run(); _set_run(r2, "   " + d, SANS, 11, INK_70)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE — CLOSING / CONTACTS
# ═══════════════════════════════════════════════════════════════════════════
def s_closing():
    s = slide(); bg(s, NIGHTBLUE)
    rect(s, 0, Inches(4.7), EMU_W, Inches(2.8), fill=DEEP_INDIGO)
    line_shape(s, Inches(0.9), Inches(1.5), Inches(1.7), Pt(3), EMBER)
    if os.path.exists(LOGO_WHITE):
        s.shapes.add_picture(LOGO_WHITE, Inches(0.85), Inches(0.7), height=Inches(0.5))
    text(s, Inches(0.9), Inches(1.75), Inches(11.5), Inches(1.6),
         [[("Thank you — over to you", SERIF, 40, WHITE)],
          [("The app is small, self-contained and fully documented. Migrating it is a "
            "provision-and-deploy exercise, not a rewrite.", SERIF, 18, CREAM)]],
         line_spacing=1.05, space_after=8)
    text(s, Inches(0.9), Inches(4.15), Inches(11), Inches(0.4),
         [[("QUESTIONS DURING TRANSITION", SANS_SEMI, 10.5, SANDSTONE, True, False, 2.0)]])
    # contacts
    conts = [
        ("Handover from", "Aleksander Biesaga", "aleksander.biesaga@teroxx.com"),
        ("Research owners", "Jannick Broering · Leonardo Larieira", "product & daily use"),
        ("IT owner", "Emil Jorgensen", "platform & deployment"),
    ]
    x = 0.9; w = 3.9
    for lab, name, sub in conts:
        text(s, Inches(x), Inches(5.0), Inches(w - 0.2), Inches(1.3),
             [[(lab, SANS_SEMI, 10.5, EMBER, True, False, 1.0)],
              [(name, SANS_SEMI, 14, WHITE, True)],
              [(sub, SANS, 11, CREAM)]], line_spacing=1.12, space_after=4)
        x += w
    text(s, Inches(0.9), Inches(6.95), Inches(11), Inches(0.3),
         [[("Teroxx  ·  Portfolio Allocation App  ·  Handover  ·  Confidential", SANS, 9, INK_40)]])


# ── build ───────────────────────────────────────────────────────────────────
s_cover()
s_contents()
s_divider("1", "What the app does", ["Capabilities & the advisor workflow", "The investable universes & baskets", "The allocation model", "Where the data comes from"], chip="research")
s_exec()
s_tabs()
s_workflow()
s_universes()
s_model()
s_datasources()
s_divider("2", "Project status", ["What is live and working", "Known limitations to watch", "Where it runs today"], chip="both")
s_status()
s_wherelives()
s_divider("3", "Architecture", ["The stack & how it fits together", "Data, security & the proposal engine", "The environment-variable contract"], chip="it")
s_architecture()
s_data()
s_security()
s_proposal_engine()
s_envvars()
s_divider("4", "Standing it up on your cloud", ["The migration in one picture", "A step-by-step runbook", "Hosting options & credentials"], chip="it")
s_migration_overview()
s_runbook()
s_hosting_options()
s_credentials()
s_divider("5", "Ownership & handover", ["Who owns what", "First-week checklist", "Contacts & references"], chip="both")
s_ownership()
s_checklist()
s_closing()

OUT = os.path.join(HERE, "Teroxx_Portfolio_App_Handover.pptx")
prs.save(OUT)
print("Saved:", OUT)
print("Slides:", len(prs.slides._sldIdLst))
